import json
import re
from pathlib import Path
from datetime import datetime, timezone

import fitz  # PyMuPDF
import pandas as pd
from docx import Document
from pptx import Presentation

BASE = Path(r"D:\climsystems_evidence_library")
SRC = Path(r"C:\Users\Yinpeng Li\ClimSystems Dropbox\Yinpeng Li\climsystems_ai\evidence_library\01_sources")
NORM = Path(r"C:\Users\Yinpeng Li\ClimSystems Dropbox\Yinpeng Li\climsystems_ai\evidence_library\02_normalised")
META_CSV = Path(r"C:\Users\Yinpeng Li\ClimSystems Dropbox\Yinpeng Li\climsystems_ai\evidence_library\05_metadata\documents.csv")
CHUNKS_DIR = Path(r"C:\Users\Yinpeng Li\CLIMsystems Dropbox\Yinpeng Li\climsystems_ai\evidence_library\03_chunks")
CHUNKS_DIR.mkdir(parents=True, exist_ok=True)
NORM.mkdir(parents=True, exist_ok=True)
META_CSV.parent.mkdir(parents=True, exist_ok=True)

DERIVED_COLS = [
    "file_ext","file_size_bytes","file_modified_utc","page_or_slide_count",
    "toc_present","detected_title","meta_title","meta_author","meta_created",
    "meta_modified","meta_producer","top_keywords","text_extract_ok","extract_error",
    "auto_doc_type","auto_perils","auto_jurisdiction","auto_topics","auto_confidence"
]

BASE_COLS = [
    "doc_id","original_filename","title","version","date","owner","doc_type",
    "scope","perils","jurisdiction","status","source"
]

def utc_iso(ts: float) -> str:
    return datetime.fromtimestamp(ts, tz=timezone.utc).isoformat()

def top_keywords(text: str, k: int = 15):
    words = re.findall(r"[A-Za-z][A-Za-z\-]{2,}", text.lower())
    stop = {"the","and","for","with","that","this","from","are","was","were","not","but","into","over","under",
            "within","where","when","what","which","will","shall","may","can","could","should","would","also",
            "than","then","their","there","these","those","your","you","our", "www", "...", "com", "http", "e-mail"}
    freq = {}
    for w in words:
        if w in stop:
            continue
        freq[w] = freq.get(w, 0) + 1
    return [w for w, _ in sorted(freq.items(), key=lambda x: x[1], reverse=True)[:k]]

def profile_pdf(file_path: Path):
    # Robust open; if it fails, caller catches and records error
    doc = fitz.open(file_path)
    meta = doc.metadata or {}
    page_count = doc.page_count

    toc = []
    try:
        toc = doc.get_toc(simple=True)  # may be empty
    except Exception:
        toc = []

    sample = ""
    for i in range(min(2, page_count)):
        sample += doc.load_page(i).get_text("text") + "\n"
    sample = sample.strip()

    detected_title = (meta.get("title") or "").strip() or (sample.splitlines()[0].strip() if sample else "")

    outline = []
    if toc:
        for level, title, page in toc[:2000]:
            outline.append({"level": level, "title": (title or "").strip(), "page": int(page)})
    else:
        # lightweight fallback outline (optional)
        outline = []

    return {
        "type": "pdf",
        "page_count": page_count,
        "toc_present": bool(toc),
        "embedded_metadata": {
            "title": (meta.get("title") or "").strip(),
            "author": (meta.get("author") or "").strip(),
            "creator": (meta.get("creator") or "").strip(),
            "producer": (meta.get("producer") or "").strip(),
            "creationDate": (meta.get("creationDate") or "").strip(),
            "modDate": (meta.get("modDate") or "").strip(),
        },
        "detected_title": detected_title,
        "outline": outline,
        "sample_text": sample[:4000],
        "top_keywords": top_keywords(sample)
    }


def profile_docx(file_path: Path):
    d = Document(file_path)

    outline = []
    sample_paras = []
    full_text_parts = []

    for p in d.paragraphs:
        txt = (p.text or "").strip()
        if not txt:
            continue
        full_text_parts.append(txt)

        style = (p.style.name or "").lower() if p.style else ""
        if style.startswith("heading"):
            m = re.search(r"heading\s*(\d+)", style)
            level = int(m.group(1)) if m else 1
            outline.append({"level": level, "title": txt})

        if len(sample_paras) < 15:
            sample_paras.append(txt)

    sample = "\n".join(sample_paras).strip()
    full_text = "\n".join(full_text_parts).strip()

    detected_title = outline[0]["title"] if outline else (sample_paras[0] if sample_paras else "")

    core = d.core_properties
    embedded = {
        "title": (core.title or "").strip(),
        "author": (core.author or "").strip(),
        "created": core.created.isoformat() if core.created else "",
        "modified": core.modified.isoformat() if core.modified else "",
        "last_modified_by": (core.last_modified_by or "").strip(),
    }

    return {
        "type": "docx",
        "toc_present": bool(outline),
        "embedded_metadata": embedded,
        "detected_title": embedded["title"] or detected_title,
        "outline": outline[:2000],
        "sample_text": sample[:4000],
        "top_keywords": top_keywords(sample if sample else full_text),
        "paragraph_count": len(d.paragraphs),
    }

def profile_pptx(file_path: Path):
    prs = Presentation(file_path)

    outline = []
    sample_texts = []
    full_text_parts = []

    for i, slide in enumerate(prs.slides):
        title = ""
        if slide.shapes.title and slide.shapes.title.has_text_frame:
            title = (slide.shapes.title.text or "").strip()

        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text and shape.text.strip():
                texts.append(shape.text.strip())
        slide_text = "\n".join(texts).strip()
        if slide_text:
            full_text_parts.append(slide_text)

        outline.append({"slide": i + 1, "title": title})

        if len(sample_texts) < 8 and slide_text:
            sample_texts.append(slide_text)

    sample = "\n\n".join(sample_texts).strip()
    full_text = "\n\n".join(full_text_parts).strip()

    detected_title = ""
    for item in outline:
        if item.get("title"):
            detected_title = item["title"]
            break

    return {
        "type": "pptx",
        "toc_present": True,
        "embedded_metadata": {},
        "detected_title": detected_title,
        "outline": outline[:2000],
        "sample_text": sample[:4000],
        "top_keywords": top_keywords(sample if sample else full_text),
        "slide_count": len(prs.slides),
    }

def extract_text_pdf(file_path: Path):
    """Return list of blocks: [{page, text}]"""
    doc = fitz.open(file_path)
    blocks = []
    for i in range(doc.page_count):
        txt = doc.load_page(i).get_text("text").strip()
        if txt:
            blocks.append({"page": i + 1, "text": txt})
    return blocks

def extract_text_docx(file_path: Path):
    """Return list of blocks: [{para_index, text}]"""
    d = Document(file_path)
    blocks = []
    pi = 0
    for p in d.paragraphs:
        txt = (p.text or "").strip()
        if not txt:
            continue
        blocks.append({"para": pi, "text": txt})
        pi += 1
    return blocks

def extract_text_pptx(file_path: Path):
    """Return list of blocks: [{slide, text}]"""
    prs = Presentation(file_path)
    blocks = []
    for i, slide in enumerate(prs.slides):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text and shape.text.strip():
                texts.append(shape.text.strip())
        txt = "\n".join(texts).strip()
        if txt:
            blocks.append({"slide": i + 1, "text": txt})
    return blocks

def normalize_document(file_path: Path):
    """Returns (doc_type, blocks)"""
    ext = file_path.suffix.lower()
    if ext == ".pdf":
        return "pdf", extract_text_pdf(file_path)
    if ext == ".docx":
        return "docx", extract_text_docx(file_path)
    if ext == ".pptx":
        return "pptx", extract_text_pptx(file_path)
    return "unknown", []

def write_text_json(outdir: Path, doc_id: str, file_type: str, blocks: list):
    payload = {
        "doc_id": doc_id,
        "type": file_type,
        "generated_at_utc": datetime.now(timezone.utc).isoformat(),
        "blocks": blocks,
    }
    (outdir / "text.json").write_text(json.dumps(payload, indent=2), encoding="utf-8")

DOC_TYPE_KEYWORDS = {
    "methodology": ["methodology", "method", "approach", "calibration", "validation", "technical note"],
    "dictionary": ["dictionary", "definitions", "glossary", "indicator", "variable list"],
    "standard": ["issb", "ifrs", "aasb", "asrs", "tcfd", "disclosure", "standard"],
    "report": ["report", "assessment", "results", "executive summary"],
    "presentation": ["agenda", "slide", "presentation"],
    "qa": ["qa", "quality assurance", "verification", "testing"],
}

PERIL_KEYWORDS = {
    "flood": ["flood", "inundation", "riverine", "pluvial", "storm surge"],
    "heat": ["heat", "temperature", "extreme heat"],
    "fire": ["wildfire", "bushfire", "fire weather"],
    "wind": ["wind", "cyclone", "storm"],
    "slr": ["sea level", "sea-level", "slr", "coastal"],
    "drought": ["drought", "water stress"],
}

JURIS_KEYWORDS = {
    "nz": ["new zealand", "nz", "aotearoa", "mf e", "ministry for the environment", "mbie", "stats nz", "environment canterbury", "ecan", "auckland council", "wellington council", "otago regional council", "waikato regional council","resource management act", "rma","national policy statement", "nes-f", "national environmental standards", "niwa", "metservice","lawa",],
    "au": ["australia", "australian", "au","apra", "asic","bom", "bureau of meteorology","csiro","geoscience australia","nsw", "new south wales","victoria", "queensland", "qld","south australia", "sa","western australia", "wa","tasmania", "tas","commonwealth of australia",],
    "uk": ["united kingdom", "uk", "britain", "great britain","england", "scotland", "wales", "northern ireland","environment agency","met office","defra","hm treasury","financial conduct authority", "fca",],
    "eu": ["european union", "eu","european commission","european parliament","european environment agency", "eea","csrd","taxonomy regulation","sfdr","eurostat",],
    "us": ["united states", "usa", "u.s.", "us ","epa", "environmental protection agency","noaa","usgs","sec","california", "caltrans","fema","national weather service",],
    "apac": ["asia-pacific", "apac", "asean"],
    "global": ["global","international","worldwide","multinational","cross-border","issb","ifrs","tcfd","tnfd","ipcc","world bank","imf","oecd","unep",]
}

def _score_keywords(text: str, keywords: list[str]) -> int:
    t = text.lower()
    return sum(1 for k in keywords if k in t)

def classify_document(filename: str, blocks: list[dict]) -> dict:
    # use first ~20 blocks for speed
    sample_text = " ".join(b["text"] for b in blocks[:20]).lower()
    name_text = filename.lower()
    combined = f"{name_text} {sample_text}"

    # doc type
    best_type, best_score = "misc", 0
    for dt, kws in DOC_TYPE_KEYWORDS.items():
        s = _score_keywords(combined, kws)
        if s > best_score:
            best_type, best_score = dt, s

    # perils (multi allowed)
    peril_hits = []
    for p, kws in PERIL_KEYWORDS.items():
        if _score_keywords(combined, kws) > 0:
            peril_hits.append(p)
    auto_perils = ";".join(peril_hits) if peril_hits else ""

    # jurisdiction
    juris = ""
    for j, kws in JURIS_KEYWORDS.items():
        if _score_keywords(combined, kws) > 0:
            juris = j
            break

    # simple confidence
    conf = min(0.95, 0.20 + 0.10 * best_score)
    return {
        "auto_doc_type": best_type,
        "auto_perils": auto_perils,
        "auto_jurisdiction": juris,
        "auto_topics": "",  # you can add later
        "auto_confidence": f"{conf:.2f}",
    }

def load_or_init_documents_csv():
    if META_CSV.exists() and META_CSV.stat().st_size > 0:
        df = pd.read_csv(META_CSV, dtype=str, keep_default_na=False)
        df = df.astype("object")
    else:
        df = pd.DataFrame(columns=BASE_COLS + DERIVED_COLS)

    # Ensure required columns exist
    for c in BASE_COLS + DERIVED_COLS:
        if c not in df.columns:
            df[c] = ""
    return df

def upsert_doc_row(df: pd.DataFrame, doc_id: str, original_filename: str, source: str):
    hit = df.index[df["doc_id"] == doc_id]
    if len(hit) == 0:
        # Auto-register row (draft)
        row = {c: "" for c in df.columns}
        row.update({
            "doc_id": doc_id,
            "original_filename": original_filename,
            "status": "draft",
            "source": source
        })
        df = pd.concat([df, pd.DataFrame([row])], ignore_index=True)
    return df

def set_derived_fields(df: pd.DataFrame, doc_id: str, derived: dict):
    i = df.index[df["doc_id"] == doc_id][0]
    for k, v in derived.items():
        #print(k,v)
        if k in df.columns:
            #print("yes")
            df.loc[i, k] = str(v)
            #print(df.loc[i, k])
    return df

def strip_surrogates(x):
    if isinstance(x, str):
        return x.encode("utf-8", errors = "ignore").decode("utf-8")
    return x


def chunk_text_blocks(
    doc_id: str,
    blocks: list[dict],
    max_chars: int = 2500,
    overlap_chars: int = 250,
):
    """
    Produce chunks with stable chunk_id and optional source anchors.
    Uses char-length heuristic, with WORD-SAFE overlap to avoid cutting words.
    """
    chunks = []
    chunk_id = 0
    buf = ""
    anchors = []  # collect pages/slides/paras included

    def safe_overlap_tail(prev_text: str, n: int) -> str:
        """Take last n chars, but drop any partial leading word so we start on a word boundary."""
        if n <= 0 or not prev_text:
            return ""
        tail = prev_text[-n:]
        # If tail begins mid-word, drop that partial token.
        # Example: "... vulnerab" -> drop "vulnerab" so next chunk starts cleanly.
        tail = re.sub(r"^\S+\s*", "", tail)
        return tail.strip()

    def flush():
        nonlocal chunk_id, buf, anchors
        text = buf.strip()
        if text:
            chunks.append({
                "doc_id": doc_id,
                "chunk_id": chunk_id,
                "citation": f"DOC:{doc_id}.c{chunk_id}",
                "text": text,
                "anchors": anchors[:]  # e.g., [{"page": 3}, {"page": 4}]
            })
            chunk_id += 1
        buf = ""
        anchors = []

    for b in blocks:
        t = (b.get("text") or "").strip()
        if not t:
            continue

        # add anchor info (page/slide/para)
        anchor = {}
        if "page" in b: anchor["page"] = b["page"]
        if "slide" in b: anchor["slide"] = b["slide"]
        if "para" in b: anchor["para"] = b["para"]

        # if adding would exceed, flush with overlap
        if len(buf) + len(t) + 2 > max_chars:
            flush()

            # overlap: carry last overlap_chars from previous chunk, word-safe
            if overlap_chars > 0 and chunks:
                tail = safe_overlap_tail(chunks[-1]["text"], overlap_chars)
                if tail:
                    buf = tail + "\n"
                    # keep a small anchor history (approximate, but useful)
                    anchors = chunks[-1].get("anchors", [])[-3:]

        buf += t + "\n"
        anchors.append(anchor)

    flush()
    return chunks


def write_chunks_jsonl(doc_id: str, chunks: list[dict]):
    out_path = CHUNKS_DIR / f"{doc_id}.chunks.jsonl"
    with out_path.open("w", encoding="utf-8") as f:
        for c in chunks:
            f.write(json.dumps(c, ensure_ascii=False) + "\n")
    return out_path

def main():
    df = load_or_init_documents_csv()

    for f in SRC.iterdir():
        print(f.name)
        if not f.is_file():
            continue
        if f.suffix.lower() not in [".pdf", ".docx", ".pptx"]:
            continue

        doc_id = f.stem
        outdir = NORM / doc_id
        outdir.mkdir(parents=True, exist_ok=True)

        # Auto-register (this fixes your empty documents.csv)
        df = upsert_doc_row(df, doc_id=doc_id, original_filename=f.name, source="local")

        stat = f.stat()
        base_derived = {
            "file_ext": f.suffix.lower().lstrip("."),
            "file_size_bytes": stat.st_size,
            "file_modified_utc": utc_iso(stat.st_mtime),
            "text_extract_ok": "true", ##MAKE WORK
            "extract_error": ""
        }

        try:
            file_type, blocks = normalize_document(f)
            write_text_json(outdir, doc_id, file_type, blocks)
        except Exception as e:
            base_derived["text_extract_ok"] = "false"
            base_derived["extract_error"] = str(e)
            df = set_derived_fields(df, doc_id, base_derived)
            continue

        profile = {"doc_id": doc_id, "file_name": f.name, "generated_at_utc": datetime.now(timezone.utc).isoformat()}

        try:
            
            if f.suffix.lower() == ".pdf":
                p = profile_pdf(f)
                profile.update(p)
                base_derived["title"] = p.get("detected_title","") or meta.get("title", "")
                base_derived["page_or_slide_count"] = str(p.get("page_count", ""))
                base_derived["toc_present"] = str(bool(p.get("toc_present", False))).lower()
                base_derived["detected_title"] = p.get("detected_title", "")

                meta = p.get("embedded_metadata", {}) or {}
                base_derived["meta_title"] = meta.get("title", "")
                base_derived["meta_author"] = meta.get("author", "")
                base_derived["meta_created"] = meta.get("creationDate", "")
                base_derived["meta_modified"] = meta.get("modDate", "")
                base_derived["meta_producer"] = meta.get("producer", meta.get("creator", ""))
                base_derived["top_keywords"] = ";".join((p.get("top_keywords", []) or [])[:20])

            elif f.suffix.lower() == ".docx":
                p = profile_docx(f)
                profile.update(p)
                base_derived["title"] = p.get("detected_title","") or meta.get("title", "")
                base_derived["page_or_slide_count"] = str(p.get("paragraph_count", ""))
                base_derived["toc_present"] = str(bool(p.get("toc_present", False))).lower()
                base_derived["detected_title"] = p.get("detected_title", "")

                meta = p.get("embedded_metadata", {}) or {}
                base_derived["meta_title"] = meta.get("title", "")
                base_derived["meta_author"] = meta.get("author", "")
                base_derived["meta_created"] = meta.get("created", "")
                base_derived["meta_modified"] = meta.get("modified", "")
                base_derived["meta_producer"] = meta.get("last_modified_by", "")
                base_derived["top_keywords"] = ";".join((p.get("top_keywords", []) or [])[:20])

            elif f.suffix.lower() == ".pptx":
                p = profile_pptx(f)
                profile.update(p)

                base_derived["page_or_slide_count"] = str(p.get("slide_count", ""))
                base_derived["toc_present"] = "true"
                base_derived["detected_title"] = p.get("detected_title", "")
                base_derived["top_keywords"] = ";".join((p.get("top_keywords", []) or [])[:20])

            else:
                profile["type"] = "unknown"

        except Exception as e:
            base_derived["text_extract_ok"] = "false"
            base_derived["extract_error"] = str(e)
            profile["error"] = str(e)


                # Write profile even if extract fails
        (outdir / "profile.json").write_text(json.dumps(profile, indent=2), encoding="utf-8")
        structure = {"doc_id": profile["doc_id"], "outline": profile.get("outline", [])}
        (outdir / "structure.json").write_text(json.dumps(structure, indent=2), encoding="utf-8")

        auto = classify_document(f.name, blocks)
        for col in auto.keys():
            if col not in df.columns:
                df[col] = ""
        # Update CSV
        df = set_derived_fields(df, doc_id, base_derived)
        df = set_derived_fields(df, doc_id, auto)
        chunks = chunk_text_blocks(doc_id, blocks, max_chars=2500, overlap_chars=250)
        chunk_path = write_chunks_jsonl(doc_id, chunks)
        #print(df["auto_perils"])
        
        
    
    df = df.map(strip_surrogates)
    #print(df)
    #print(df["auto_perils"])
    df.to_csv(META_CSV, index=False, encoding="utf-8")
    print(f"Done. Updated: {META_CSV}")

if __name__ == "__main__":
    main()